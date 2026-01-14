import os
import json
import logging
import re
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE

# --- Basic Setup ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# --- Helper function to find specific placeholders ---
def _get_placeholder(slide, ph_type):
    for shape in slide.placeholders:
        if shape.placeholder_format.type == ph_type:
            return shape
    return None


# --- NEW: Markdown Paragraph Renderer ---
def _render_markdown_paragraph(paragraph, markdown_text):
    """
    Parses a simple Markdown string and adds formatted runs to a paragraph.
    Supports **bold** and *italic*.
    """
    # This regex splits the text by the Markdown markers but keeps them in the list
    segments = re.split(r'(\*\*.*?\*\*|\*.*?\*)', markdown_text)
    
    for segment in segments:
        if not segment: continue # Skip empty strings
        
        run = paragraph.add_run()
        if segment.startswith('**') and segment.endswith('**'):
            run.text = segment[2:-2]  # Remove the ** markers
            run.font.bold = True
        elif segment.startswith('*') and segment.endswith('*'):
            run.text = segment[1:-1]  # Remove the * markers
            run.font.italic = True
        else:
            run.text = segment

# --- IMPROVEMENT: Combined Smart-Fit and Markdown Rendering ---
# --- IMPROVEMENT: The rendering function now gives the subtitle special formatting ---
def _render_text_with_smart_fit_and_markdown(text_frame, data, subtitle=None):
    """
    Renders complex list data, parsing Markdown for formatting.
    If a subtitle is provided, it is rendered as a visually distinct sub-header.
    """
    text_frame.clear()
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE # Fallback

    # --- 1. RENDER THE SUBTITLE (if it exists) ---
    if subtitle:
        p_subtitle = text_frame.add_paragraph()
        p_subtitle.text = subtitle
        p_subtitle.level = 0  # Level 0 = no bullet point
        p_subtitle.font.bold = True
        # Make the subtitle font noticeably larger than the body text
        p_subtitle.font.size = Pt(24) 
        # Add space after the subtitle for visual separation
        p_subtitle.space_after = Pt(12) 

    # --- 2. RENDER THE BULLET POINTS ---
    # Calculate font size based on the main data only
    json_string = json.dumps(data)
    char_count = len(json_string)
    font_tiers = [(1200, Pt(16)), (750, Pt(18)), (450, Pt(20)), (0, Pt(22))]
    body_font_size = next((size for threshold, size in font_tiers if char_count >= threshold), Pt(22))

    def add_points_recursive(points, level):
        for point in points:
            p = text_frame.add_paragraph()
            # Start bullet points at level 1 if there was a subtitle, otherwise start at 0
            p.level = level + 1 if subtitle else level
            p.font.size = body_font_size

            text_to_render = point.get('text', '') if isinstance(point, dict) else str(point)
            _render_markdown_paragraph(p, text_to_render)

            if isinstance(point, dict) and 'children' in point:
                add_points_recursive(point['children'], level + 1)

    if isinstance(data, list):
        add_points_recursive(data, 0)
        
    # Cleanup initial empty paragraph
    if text_frame.paragraphs and not text_frame.paragraphs[0].text:
        p_to_remove = text_frame.paragraphs[0]
        p_to_remove._p.getparent().remove(p_to_remove._p)





# --- Rendering Helper Functions ---
def _render_bullet_points(text_frame, data):
    text_frame.clear(); text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    def add_points(points, level):
        for point in points:
            if isinstance(point, dict):
                p = text_frame.add_paragraph(); p.text = point.get('text', ''); p.level = level + 1
                if 'children' in point: add_points(point['children'], level + 1)
            else:
                p = text_frame.add_paragraph(); p.text = str(point); p.level = level+1
    if isinstance(data, list): add_points(data, 0)
    if text_frame.paragraphs and not text_frame.paragraphs[0].text:
        p = text_frame.paragraphs[0]; p._p.getparent().remove(p._p)

def _render_table(slide, placeholder, data):
    headers, rows_data = data.get('headers', []), data.get('rows', [])
    if not headers or not rows_data: return
    table_shape = slide.shapes.add_table(len(rows_data) + 1, len(headers), placeholder.left, placeholder.top, placeholder.width, placeholder.height)
    table = table_shape.table
    for i, header in enumerate(headers):
        cell = table.cell(0, i); cell.text = header; cell.text_frame.paragraphs[0].font.bold = True
    for r_idx, row_data in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row_data):
            table.cell(r_idx + 1, c_idx).text = str(cell_text)
    sp = placeholder.element; sp.getparent().remove(sp)

def _render_mcq(placeholder, data, answer_placeholder=None):
    """
    Renders the MCQ question into the main placeholder and the detailed
    answer with explanation into the dedicated answer_placeholder.
    """
    # --- Part 1: Render the Question and Options ---
    tf = placeholder.text_frame
    tf.clear()
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    # Add the question text
    p_question = tf.add_paragraph()
    p_question.text = data.get('question_text', '')
    p_question.font.bold = True

    # Add the options
    for option in data.get('options', []):
        p_option = tf.add_paragraph()
        p_option.text = f"{option.get('label', '')}) {option.get('text', '')}"
        p_option.level = 1

    # --- Part 2: Render the Detailed Answer and Explanation ---
    if answer_placeholder:
        answer_info = data.get('correct_answer', {})
        answer_label = answer_info.get('label', 'N/A')
        answer_explanation = answer_info.get('explanation', 'No explanation provided.')
        
        # Get the text frame of the answer placeholder
        answer_tf = answer_placeholder.text_frame
        answer_tf.clear()
        
        # Add "Correct Answer: [Label]" in bold
        p_answer_label = answer_tf
        p_answer_label.text = f" Answer: {answer_label} - {answer_explanation}"
        #p_answer_label.font.bold = True
        
        # # Add the explanation text on the next line
        # p_answer_explanation = answer_tf.add_paragraph()
        # p_answer_explanation.text = answer_explanation

def _render_matching(placeholders, data):
    """Renders a matching activity into a pair of placeholders."""
    if len(placeholders) < 2: return
    terms_ph, defs_ph = placeholders[0], placeholders[1]
    terms_tf, defs_tf = terms_ph.text_frame, defs_ph.text_frame
    terms_tf.clear(); defs_tf.clear()
    pairs = data.get('pairs', [])
    if not pairs: return
    terms = [f"{i+1}. {p['term']}" for i, p in enumerate(pairs)]
    definitions = [p['definition'] for p in pairs]
    random.shuffle(definitions)
    for term in terms: terms_tf.add_paragraph().text = term
    for i, definition in enumerate(definitions): defs_tf.add_paragraph().text = f"{chr(65+i)}. {definition}"
    
    

        


    

class PresentationAgent:
    def __init__(self, template_path: str, layout_config_path: str):
        self.template_path = template_path
        self.prs_for_layouts = Presentation(template_path)
        with open(layout_config_path, 'r', encoding='utf-8') as f:
            self.user_selections = json.load(f)['user_selections']
        logger.info(f"Agent initialized with template '{os.path.basename(template_path)}'")

    def _get_layout(self, layout_key: str):
        selection = self.user_selections.get(layout_key, self.user_selections['Content'])
        return self.prs_for_layouts.slide_layouts[selection['selected_layout_index']]

    def _collect_and_sort_slides(self, plan_data: dict) -> list:
        all_slide_items = []

        # This is the same recursive helper function, it is correct.
        def _recursive_harvester(node: dict, parent_title: str = None):
            is_child_node = bool(parent_title)
            for slide in node.get('slides', []):
                slide['parent_title'] = parent_title
                slide['is_child'] = is_child_node
                all_slide_items.append({'item_type': 'Content', 'data': slide})
            
            for child in node.get('children', []):
                # Pass the current node's title as the parent title for the next level
                _recursive_harvester(child, parent_title=node.get('title'))
            
            # This part correctly harvests activities from ANY node (parent or child)
            if 'interactive_activity' in node and 'llm_generated_content' in node['interactive_activity']:
                activity = node['interactive_activity']
                activity['parent_title'] = parent_title
                # IMPORTANT: The 'is_child' status depends on whether a parent_title was passed in.
                activity['is_child'] = is_child_node
                # We add the main node's title to the activity data so it can be used if needed
                activity['main_topic_title'] = node.get('title')
                all_slide_items.append({'item_type': 'Content', 'data': activity})

        # This is the main loop that processes the plan
        for deck in plan_data.get('deck_plans', []):
            # Framework slides are added as before
            framework_start = ([{'item_type': 'Title', 'data': s} for s in deck['sections'] if s['section_type'] == 'Title'] +
                            [{'item_type': 'Agenda', 'data': s} for s in deck['sections'] if s['section_type'] == 'Agenda'])
            all_slide_items.extend(framework_start)
            
            # --- MODIFICATION IS HERE ---
            # We now loop through the content blocks and call the harvester on each one.
            content_sections = [s for s in deck['sections'] if s['section_type'] == 'Content']
            if content_sections:
                for block in content_sections[0].get('content_blocks', []):
                    # We start the recursion for each main content block
                    # We pass the block's own title as the parent_title for its direct activities
                    _recursive_harvester(block, parent_title=None) # Start with no parent

            framework_end = ([{'item_type': 'Summary', 'data': s} for s in deck['sections'] if s['section_type'] == 'Summary'] +
                            [{'item_type': 'End', 'data': s} for s in deck['sections'] if s['section_type'] == 'End'])
            all_slide_items.extend(framework_end)
        
        return sorted(all_slide_items, key=lambda x: x['data'].get('seq_id', 999))
        
    def _determine_layout_key(self, slide_data: dict) -> str:
        # This logic is correct
        llm_content = slide_data.get('llm_generated_content', {})
        objects = llm_content.get('objects', [])
        if not objects: return "Content"
        first_obj_type = objects[0].get('content_type')
        if first_obj_type in ['multiple_choice_question', 'matching_activity']:
            return "Application_Two_Column" if first_obj_type == 'matching_activity' else "Application"
        is_child = slide_data.get('is_child', False)
        if len(objects) >= 2:
            return "Content_Two_Column_child" if is_child else "Content_Two_Column"
        return "Content_child" if is_child else "Content"

    def create_presentation_from_plan(self, plan_json_path: str, output_dir: str):
        # This orchestrator is correct
        logger.info(f"Creating presentation from plan: '{os.path.basename(plan_json_path)}'")
        with open(plan_json_path, 'r', encoding='utf-8') as f: plan_data = json.load(f)
        final_slide_list = self._collect_and_sort_slides(plan_data)
        prs = Presentation(self.template_path)
        while len(prs.slides):
            rId = prs.slides._sldIdLst[0].rId; prs.part.drop_rel(rId); del prs.slides._sldIdLst[0]
        logger.info(f"Collected a total of {len(final_slide_list)} slides to render.")
        for item in final_slide_list:
            self._add_slide_for_item(prs, item)
        unit_code = plan_data.get('deck_plans', [{}])[0].get('sections', [{}])[0].get('content', {}).get('unit_code', 'UNIT')
        week_num = plan_data.get('week', 'X')
        deck_num = plan_data.get('deck_plans', [{}])[0].get('deck_number', 'Y')
        output_filename = f"{unit_code}_Week{week_num}_Deck{deck_num}_Presentation.pptx"
        output_path = os.path.join(output_dir, output_filename)
        prs.save(output_path)
        logger.info(f"Successfully created presentation: '{output_path}'")

    def _add_slide_for_item(self, prs: Presentation, item: dict):
        item_type = item['item_type']
        item_data = item['data']
        layout_key = self._determine_layout_key(item_data) if item_type == 'Content' else item_type
        slide = prs.slides.add_slide(self._get_layout(layout_key))
        render_map = {
            'Title': self._render_title_slide, 'Agenda': self._render_agenda_slide,
            'Summary': self._render_summary_slide, 'End': self._render_end_slide,
            'Divider': self._render_divider_slide, 'Content': self._render_content_slide
        }
        if item_type in render_map: render_map[item_type](slide, item_data)
    
    def _render_title_slide(self, slide, data):
        content = data.get('content', {})
        title_ph = _get_placeholder(slide, PP_PLACEHOLDER_TYPE.CENTER_TITLE)
        if title_ph: title_ph.text = content.get('week_topic', '')
        subtitle_ph = _get_placeholder(slide, PP_PLACEHOLDER_TYPE.SUBTITLE)
        if subtitle_ph: subtitle_ph.text = f"{content.get('deck_title', '')}\n{content.get('unit_name', '')} | {content.get('unit_code', '')}"

    def _render_agenda_slide(self, slide, data):
        content = data.get('content', {})
        title_ph = _get_placeholder(slide, PP_PLACEHOLDER_TYPE.TITLE)
        if title_ph: title_ph.text = content.get('title', 'Agenda')
        body_ph = _get_placeholder(slide, PP_PLACEHOLDER_TYPE.OBJECT)
        if body_ph: _render_bullet_points(body_ph.text_frame, content.get('items', []))

   

    def _render_divider_slide(self, slide, data):
        content = data.get('content', {})
        title_ph = _get_placeholder(slide, PP_PLACEHOLDER_TYPE.TITLE)
        if title_ph: title_ph.text = content.get('title', '')

    def _render_end_slide(self, slide, data):
        content = data.get('content', {})
        title_ph = _get_placeholder(slide, PP_PLACEHOLDER_TYPE.TITLE)
        if title_ph: title_ph.text = content.get('text', 'Thank You & Questions?')
        
    def _render_summary_slide(self, slide, data):
        title_ph = _get_placeholder(slide, PP_PLACEHOLDER_TYPE.TITLE)
        if title_ph: title_ph.text = "Summary & Key Takeaways"
        body_ph = _get_placeholder(slide, PP_PLACEHOLDER_TYPE.OBJECT)
        if body_ph and 'llm_generated_content' in data:
            summary_data = data['llm_generated_content'].get('objects', [{}])[0].get('data', [])
            _render_bullet_points(body_ph.text_frame, summary_data)

    # <<< THIS IS THE CORRECTED METHOD >>>
    def _render_content_slide(self, slide, data):
        content = data.get('llm_generated_content', {})
        if not content: return
        layout_key = self._determine_layout_key(data)
        title_ph = _get_placeholder(slide, PP_PLACEHOLDER_TYPE.TITLE)
        subtitle_ph = _get_placeholder(slide, PP_PLACEHOLDER_TYPE.SUBTITLE)
        answer_placeholder = _get_placeholder(slide, PP_PLACEHOLDER_TYPE.BODY)
        

        # --- 2. Render Titles Based on Layout Key ---
        if "Application" in layout_key:
            if title_ph: title_ph.text = content.get('subtitle', 'Knowledge Check')
            if subtitle_ph: subtitle_ph.text = content.get('title', "Let's Apply This!")
        elif layout_key == "Content_child" or layout_key == "Content_Two_Column_child":
            if title_ph: title_ph.text = data.get('parent_title', '')
            if subtitle_ph: subtitle_ph.text = content.get('title', '')
        else: # Standard "Content" or "Content_Two_Column"
            if title_ph: title_ph.text = content.get('title', '')
            # For the main parent layouts, the subtitle comes from the JSON
            if subtitle_ph: subtitle_ph.text = content.get('subtitle', '')
            
            
        # Resize    
        if title_ph: title_ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if subtitle_ph: subtitle_ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        body_placeholders = sorted([p for p in slide.placeholders if p.placeholder_format.type == PP_PLACEHOLDER_TYPE.OBJECT], key=lambda p: p.left)
        
        objects = content.get('objects', [])

        
        if "Two_Column" in layout_key and len(objects) >= 2:
            # For two-column layouts, we look for a BODY placeholder for the subtitle
            body_subtitle_ph = _get_placeholder(slide, PP_PLACEHOLDER_TYPE.BODY)
            if body_subtitle_ph and content.get('subtitle'):
                p = body_subtitle_ph.text_frame
                p.text = content.get('subtitle')
                
            
            # Render the first object in the first placeholder
            obj1_content_type, obj1_data = objects[0].get('content_type'), objects[0].get('data')
            if obj1_content_type == 'bullet_points':
                _render_text_with_smart_fit_and_markdown(body_placeholders[0].text_frame, obj1_data)
            
            # Render the second object in the second placeholder
            obj2_content_type, obj2_data = objects[1].get('content_type'), objects[1].get('data')
            if obj2_content_type == 'bullet_points':
                _render_text_with_smart_fit_and_markdown(body_placeholders[1].text_frame, obj2_data)
        
        
        else:
            for i, obj in enumerate(objects):
                if i >= len(body_placeholders): break
                placeholder = body_placeholders[i]
                content_type, obj_data = obj.get('content_type'), obj.get('data')
                object_subtitle = content.get('subtitle', None) 

                # --- IMPROVEMENT: Call the new, more capable rendering functions ---
                if content_type in ['bullet_points', 'description', 'explanation', 'timeline', 'process', 'cycle', 'hierarchy', 'case_study', 'example']:
                    _render_text_with_smart_fit_and_markdown(placeholder.text_frame, obj_data, subtitle=object_subtitle)
                    
                elif content_type == 'table':
                    _render_table(slide, placeholder, obj_data)
                elif content_type == 'multiple_choice_question':
                    _render_mcq(placeholder, obj_data, answer_placeholder)
                    if slide.has_notes_slide:
                        answer_info = obj_data.get('correct_answer', {})
                        slide.notes_slide.notes_text_frame.text = f"Answer: {answer_info.get('label', '')}. {answer_info.get('explanation', '')}"
                elif content_type == 'matching_activity':
                    # _render_matching is not defined in the provided snippet but would be called here
                    pass
                else: 
                    # Fallback to the simple bullet point renderer if type is unknown
                    _render_bullet_points(placeholder.text_frame, obj_data)